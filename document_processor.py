"""
Document processing module for handling various file formats.
Supports PDF, DOCX, PPTX, EPUB, and video/audio files.
"""

import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
import mimetypes

# Document processing libraries (with optional imports)
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PyPDF2 = None
    PYPDF2_AVAILABLE = False

# Image processing libraries
try:
    import fitz  # PyMuPDF for image extraction
    PYMUPDF_AVAILABLE = True
except ImportError:
    fitz = None
    PYMUPDF_AVAILABLE = False

try:
    from PIL import Image
    import io
    import base64
    PIL_AVAILABLE = True
except ImportError:
    Image = None
    PIL_AVAILABLE = False

# OpenAI for multimodal image analysis (lazy loading)
try:
    from openai import OpenAI
    import os
    OPENAI_AVAILABLE = True
    openai_client = None  # Initialize later when needed
except ImportError:
    OpenAI = None
    openai_client = None
    OPENAI_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    docx = None
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False

try:
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    EPUB_AVAILABLE = True
except ImportError:
    ebooklib = None
    epub = None
    BeautifulSoup = None
    EPUB_AVAILABLE = False

# Video/Audio processing (optional)
try:
    import whisper
    WHISPER_AVAILABLE = True
except ImportError:
    whisper = None
    WHISPER_AVAILABLE = False

try:
    import ffmpeg
    FFMPEG_AVAILABLE = True
except ImportError:
    ffmpeg = None
    FFMPEG_AVAILABLE = False

from config import Config

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Process various document formats and extract text content."""
    
    def __init__(self):
        """Initialize the document processor."""
        self.config = Config()
        self.whisper_model = None
    
    def load_whisper_model(self):
        """Load Whisper model for audio/video transcription."""
        if self.whisper_model is None:
            try:
                model_size = self.config.model_config['whisper']['model_size']
                device = self.config.model_config['whisper']['device']
                self.whisper_model = whisper.load_model(model_size, device=device)
                logger.info(f"Loaded Whisper {model_size} model on {device}")
            except Exception as e:
                logger.error(f"Failed to load Whisper model: {e}")
                raise
    
    def process_file(self, file_path: Path, is_syllabus: bool = False) -> Dict[str, Any]:
        """
        Process a file and extract its content.
        
        Args:
            file_path: Path to the file to process
            is_syllabus: Whether this file should be weighted as syllabus content
            
        Returns:
            Dictionary containing processed document information
        """
        logger.info(f"Processing file: {file_path}")
        
        # Determine file type
        mime_type, _ = mimetypes.guess_type(str(file_path))
        file_extension = file_path.suffix.lower()
        
        # Extract content based on file type
        try:
            if file_extension == '.pdf':
                if not PYPDF2_AVAILABLE:
                    raise ValueError("PyPDF2 not installed. Cannot process PDF files.")
                content = self._process_pdf(file_path)
            elif file_extension == '.docx':
                if not DOCX_AVAILABLE:
                    raise ValueError("python-docx not installed. Cannot process DOCX files.")
                content = self._process_docx(file_path)
            elif file_extension == '.pptx':
                if not PPTX_AVAILABLE:
                    raise ValueError("python-pptx not installed. Cannot process PPTX files.")
                content = self._process_pptx(file_path)
            elif file_extension == '.epub':
                if not EPUB_AVAILABLE:
                    raise ValueError("ebooklib not installed. Cannot process EPUB files.")
                content = self._process_epub(file_path)
            elif file_extension in ['.mp4', '.avi', '.mov', '.mp3', '.wav']:
                if not WHISPER_AVAILABLE:
                    raise ValueError("whisper not installed. Cannot process audio/video files.")
                content = self._process_audio_video(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")
            
            # Create document metadata
            document = {
                'file_path': str(file_path),
                'file_name': file_path.name,
                'file_type': file_extension,
                'content': content,
                'is_syllabus': is_syllabus,
                'syllabus_weight': self.config.chunk_config['syllabus_weight'] if is_syllabus else 1.0,
                'character_count': len(content),
                'word_count': len(content.split()) if content else 0,
            }
            
            logger.info(f"Successfully processed {file_path.name}: {len(content)} characters")
            return document
            
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            raise
    
    def analyze_image_with_gpt4v(self, image_data: bytes, page_num: int = None) -> str:
        """Analyze image content using GPT-4V."""
        global openai_client
        
        if not OPENAI_AVAILABLE:
            return "[Image content - multimodal analysis not available]"
        
        # Lazy initialization of OpenAI client
        if openai_client is None:
            api_key = os.environ.get("OPENAI_API_KEY")
            if not api_key:
                return "[Image content - OpenAI API key not configured]"
            try:
                openai_client = OpenAI(api_key=api_key)
            except Exception as e:
                logger.error(f"Failed to initialize OpenAI client: {e}")
                return "[Image content - OpenAI initialization failed]"
        
        try:
            # Encode image to base64
            base64_image = base64.b64encode(image_data).decode('utf-8')
            
            page_context = f" on page {page_num}" if page_num else ""
            
            response = openai_client.chat.completions.create(
                model="gpt-4o",  # Latest multimodal model
                messages=[
                    {
                        "role": "user", 
                        "content": [
                            {
                                "type": "text",
                                "text": f"Analyze this image from a business/academic document{page_context}. Describe charts, graphs, diagrams, tables, key data points, and any text visible in the image. Focus on extracting meaningful business insights and data that would be useful for learning and analysis."
                            },
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:image/png;base64,{base64_image}"}
                            }
                        ]
                    }
                ],
                max_tokens=1000
            )
            
            return f"[IMAGE ANALYSIS{page_context}]: {response.choices[0].message.content}"
            
        except Exception as e:
            logger.error(f"Error analyzing image: {e}")
            return "[Image content - analysis failed]"

    def _process_pdf(self, file_path: Path) -> str:
        """Extract text and analyze images from PDF files."""
        content_parts = []
        
        # First try PyMuPDF for comprehensive extraction
        if PYMUPDF_AVAILABLE and PIL_AVAILABLE:
            try:
                doc = fitz.open(file_path)
                
                for page_num in range(doc.page_count):
                    page = doc[page_num]
                    
                    # Extract text
                    text = page.get_text()
                    if text.strip():
                        content_parts.append(f"[Page {page_num + 1} Text]\n{text}")
                    
                    # Extract and analyze images
                    image_list = page.get_images()
                    for img_index, img in enumerate(image_list):
                        try:
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            
                            if pix.n - pix.alpha < 4:  # Valid image
                                img_data = pix.tobytes("png")
                                
                                # Analyze image with GPT-4V if available and configured
                                if OPENAI_AVAILABLE and os.environ.get("OPENAI_API_KEY"):
                                    analysis = self.analyze_image_with_gpt4v(img_data, page_num + 1)
                                    content_parts.append(analysis)
                                else:
                                    content_parts.append(f"[IMAGE on page {page_num + 1}] - Image present (configure OPENAI_API_KEY for analysis)")
                            
                            pix = None  # Clean up
                            
                        except Exception as e:
                            logger.warning(f"Error processing image {img_index} on page {page_num + 1}: {e}")
                            continue
                
                doc.close()
                return "\n\n".join(content_parts) if content_parts else ""
                
            except Exception as e:
                logger.warning(f"PyMuPDF processing failed for {file_path}: {e}, falling back to PyPDF2")
        
        # Fallback to PyPDF2 for text-only extraction
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text_content = []
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content.append(f"[Page {page_num + 1}]\n{page_text}")
                    except Exception as e:
                        logger.warning(f"Could not extract text from page {page_num + 1}: {e}")
                
                return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            raise
    
    def _process_docx(self, file_path: Path) -> str:
        """Extract text from DOCX file."""
        try:
            doc = docx.Document(file_path)
            text_content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text)
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            raise
    
    def _process_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX file."""
        try:
            presentation = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = [f"[Slide {slide_num + 1}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the slide number
                    text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
            raise
    
    def _process_epub(self, file_path: Path) -> str:
        """Extract text from EPUB file."""
        try:
            book = epub.read_epub(str(file_path))
            text_content = []
            
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    # Parse HTML content
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    
                    # Extract text, preserving some structure
                    for script in soup(["script", "style"]):
                        script.decompose()
                    
                    text = soup.get_text()
                    
                    # Clean up whitespace
                    lines = (line.strip() for line in text.splitlines())
                    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                    text = ' '.join(chunk for chunk in chunks if chunk)
                    
                    if text.strip():
                        # Try to get chapter title
                        title_tag = soup.find(['h1', 'h2', 'title'])
                        chapter_title = title_tag.get_text() if title_tag else f"Chapter {len(text_content) + 1}"
                        text_content.append(f"[{chapter_title}]\n{text}")
            
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Error processing EPUB {file_path}: {e}")
            raise
    
    def _process_audio_video(self, file_path: Path) -> str:
        """Transcribe audio/video files using Whisper."""
        try:
            # Load Whisper model if not already loaded
            if self.whisper_model is None:
                self.load_whisper_model()
            
            logger.info(f"Transcribing audio/video file: {file_path.name}")
            
            # Use Whisper to transcribe
            result = self.whisper_model.transcribe(str(file_path))
            
            # Format the transcription with timestamps if available
            if 'segments' in result:
                formatted_segments = []
                for segment in result['segments']:
                    start_time = self._format_timestamp(segment.get('start', 0))
                    text = segment.get('text', '').strip()
                    if text:
                        formatted_segments.append(f"[{start_time}] {text}")
                
                transcription = "\n".join(formatted_segments)
            else:
                transcription = result.get('text', '')
            
            logger.info(f"Transcription completed: {len(transcription)} characters")
            return transcription
            
        except Exception as e:
            logger.error(f"Error transcribing {file_path}: {e}")
            raise
    
    def _format_timestamp(self, seconds: float) -> str:
        """Format seconds into MM:SS format."""
        minutes = int(seconds // 60)
        seconds = int(seconds % 60)
        return f"{minutes:02d}:{seconds:02d}"
    
    def get_supported_formats(self) -> List[str]:
        """Return list of supported file formats."""
        return ['.pdf', '.docx', '.pptx', '.epub', '.mp4', '.avi', '.mov', '.mp3', '.wav']
    
    def is_supported_format(self, file_path: Path) -> bool:
        """Check if file format is supported."""
        return file_path.suffix.lower() in self.get_supported_formats()
    
    def is_syllabus_file(self, file_path: Path) -> bool:
        """Automatically detect if a file is likely a syllabus based on filename."""
        filename_lower = file_path.name.lower()
        syllabus_keywords = [
            'syllabus', 'outline', 'curriculum', 'overview', 
            'course_info', 'course_outline', 'program_outline',
            'schedule', 'agenda', 'course_description'
        ]
        
        return any(keyword in filename_lower for keyword in syllabus_keywords)
    
    def process_directory(self, directory_path: Path) -> List[Dict[str, Any]]:
        """
        Process all supported files in a directory recursively.
        
        Args:
            directory_path: Path to the directory to process
            
        Returns:
            List of processed document dictionaries
        """
        if not directory_path.exists():
            raise ValueError(f"Directory does not exist: {directory_path}")
        
        processed_files = []
        
        # Walk through directory recursively
        for file_path in directory_path.rglob('*'):
            if file_path.is_file() and self.is_supported_format(file_path):
                try:
                    # Auto-detect syllabus files
                    is_syllabus = self.is_syllabus_file(file_path)
                    
                    # Process the file
                    processed_doc = self.process_file(file_path, is_syllabus=is_syllabus)
                    processed_files.append(processed_doc)
                    
                    logger.info(f"Processed {'[SYLLABUS] ' if is_syllabus else ''}{file_path.name}")
                    
                except Exception as e:
                    logger.warning(f"Failed to process {file_path}: {e}")
                    continue
        
        logger.info(f"Processed {len(processed_files)} files from directory {directory_path}")
        return processed_files
